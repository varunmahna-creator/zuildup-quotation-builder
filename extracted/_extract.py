#!/usr/bin/env python3
"""ZuildUp quotation-builder source extractor.
Reads PPTs + DOCX from src_* dirs and writes structured data into extracted/
and slide images into assets/lookbook/{package}/.
"""
import os, re, json, zipfile, shutil, sys, traceback
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu
from docx import Document
from PIL import Image

ROOT = Path("/opt/ocplatform/workspace/zuildup/quotation-builder")
SRC_PPT  = ROOT/"src_ppt"
SRC_DOCX = ROOT/"src_docx"
EXTR     = ROOT/"extracted"
LOOK     = ROOT/"assets/lookbook"

PKG_MAP = {
    "Premium.pptx":      "premium",
    "Platinum_WIP.pptx": "platinum",
    "Royale.pptx":       "royale",
}
DOCX_MAP = {
    "Premium.docx":  "premium",
    "Platinum.docx": "platinum",
    "Royale.docx":   "royale",
}

BRAND_TOKENS = [
    "Jaquar","Hindware","Kohler","Cera","Parryware","Asian Paints","Berger",
    "Birla White","JK Lakshmi","UltraTech","Ambuja","JSW","Tata","Jindal",
    "Havells","Anchor","Legrand","Schneider","Crompton","Polycab","Finolex",
    "Greenply","Century","Saint-Gobain","Nippon","Dulux","Kajaria","Somany",
    "Johnson","Orient Bell","Ozone","Yale","Dorma","Fenesta","Italian",
    "Roca","Toto","Grohe","American Standard","Hettich","Hafele","Ebco",
]
CATEGORY_KW = {
    "bathroom":  ["bathroom","bath","sanitary","cp fitting","faucet","wc","wash basin","shower","geyser"],
    "kitchen":   ["kitchen","sink","chimney","hob","modular","platform"],
    "structure": ["structure","foundation","footing","beam","column","slab","plinth","rcc"],
    "cement":    ["cement","concrete","plaster","mortar"],
    "steel":     ["steel","tmt","reinforce"],
    "electrical":["electric","wiring","switch","mcb","cable","point","light fitting","fan"],
    "flooring":  ["flooring","tile","marble","granite","wooden floor","vitrified"],
    "doors_windows":["door","window","frame","aluminium","upvc","fenesta","main door"],
    "water":     ["water","tank","pump","plumbing","pipe","cpvc","upvc pipe"],
    "ceiling":   ["ceiling","false ceiling","pop","gypsum"],
    "safety":    ["safety","cctv","intercom","alarm","fire"],
    "paint":     ["paint","emulsion","primer","putty","texture"],
    "waterproofing":["waterproof","waterproofing"],
    "parapet":   ["parapet","compound wall","boundary"],
    "general":   ["general","aspect","misc","contingency","supervision"],
}

AMT_RE = re.compile(r"(?:₹|rs\.?|inr)\s*[\d,]+(?:\.\d+)?", re.I)
NUM_RE = re.compile(r"[\d,]+(?:\.\d+)?")

def looks_like_lineitem_row(cells):
    text = " | ".join(c.strip() for c in cells if c)
    has_label = any(re.search(r"[A-Za-z]{4,}", c) for c in cells)
    has_amt = AMT_RE.search(text) is not None
    if not has_amt:
        # Numeric > 100?
        for c in cells:
            for m in NUM_RE.findall(c.replace(",","")):
                try:
                    if float(m) > 100:
                        has_amt = True; break
                except: pass
            if has_amt: break
    return has_label and has_amt

def label_from_row(cells):
    # Pick the longest mostly-text cell as label
    best = ""
    for c in cells:
        cs = c.strip()
        letters = sum(ch.isalpha() for ch in cs)
        if letters >= 3 and len(cs) > len(best):
            best = cs
    return best[:200]

def amount_from_row(cells):
    text = " | ".join(cells)
    m = AMT_RE.search(text)
    if m: return m.group(0)
    # fallback: largest number
    biggest = None
    for c in cells:
        for tok in NUM_RE.findall(c.replace(",","")):
            try:
                v = float(tok)
                if v > 100 and (biggest is None or v > biggest):
                    biggest = v
            except: pass
    return f"{biggest:.0f}" if biggest else None

def brand_from_row(cells):
    text = " | ".join(cells)
    found = []
    for b in BRAND_TOKENS:
        if re.search(rf"\b{re.escape(b)}\b", text, re.I):
            found.append(b)
    return ", ".join(found) if found else None

def categorise(label):
    lab = label.lower()
    cats = []
    for cat, kws in CATEGORY_KW.items():
        if any(kw in lab for kw in kws):
            cats.append(cat)
    return cats

def extract_pptx(pptx_path: Path, pkg: str):
    import os as _os
    pptx_path = Path(str(pptx_path)).resolve()
    _abs = str(pptx_path)
    _ex_pathlib = pptx_path.exists()
    _ex_os = _os.path.exists(_abs)
    _cwd = _os.getcwd()
    print(f"[ppt] {pkg} <- {_abs}  pathlib_exists={_ex_pathlib} os_exists={_ex_os} cwd={_cwd}", flush=True)
    if not _ex_os:
        # Try opening blind to see actual error
        try:
            with open(_abs, "rb") as _f:
                print("BLIND OPEN OK, first bytes:", _f.read(4))
        except Exception as e:
            print("BLIND OPEN ERR:", e)
        # Try parent listdir
        try:
            print("PARENT listdir:", _os.listdir(_os.path.dirname(_abs)))
        except Exception as e:
            print("PARENT listdir ERR:", e)
    out_dir = EXTR/pkg; safe_mkdir(out_dir)
    look_dir = LOOK/pkg; safe_mkdir(look_dir)
    slides_data = []
    raw_text_lines = []
    raw_lineitems = []

    # First, harvest all media via the underlying zip — gives reliable extension
    media_by_filename = {}
    with zipfile.ZipFile(pptx_path) as z:
        for name in z.namelist():
            if name.startswith("ppt/media/"):
                media_by_filename[name] = z.read(name)

    pres = Presentation(str(pptx_path))
    for i, slide in enumerate(pres.slides, start=1):
        sd = {"slide_no": i, "title": None, "texts": [], "tables": [], "notes": None, "image_files": []}
        try:
            if slide.shapes.title and slide.shapes.title.text:
                sd["title"] = slide.shapes.title.text.strip()
        except Exception:
            pass
        # Text + tables
        for shape in slide.shapes:
            try:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        t = "".join(r.text for r in p.runs).strip()
                        if t: sd["texts"].append(t)
                if shape.has_table:
                    tbl = []
                    for row in shape.table.rows:
                        tbl.append([cell.text.strip() for cell in row.cells])
                    sd["tables"].append(tbl)
                # Picture shapes
                if shape.shape_type == 13 and hasattr(shape, "image"):
                    blob = shape.image.blob
                    ext = shape.image.ext or "bin"
                    fname = f"slide-{i:02d}-img-{len(sd['image_files'])+1:02d}.{ext}"
                    fp = look_dir / fname
                    fp.write_bytes(blob)
                    sd["image_files"].append(str(fp.relative_to(ROOT)))
            except Exception as e:
                print(f"  shape err slide {i}: {e}")
        # Notes
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                sd["notes"] = slide.notes_slide.notes_text_frame.text.strip() or None
        except Exception:
            pass
        # Pick "primary" image: largest by file size among this slide's image_files
        if sd["image_files"]:
            best = max(sd["image_files"], key=lambda rp: (ROOT/rp).stat().st_size)
            ext = Path(best).suffix
            primary = look_dir / f"slide-{i:02d}{ext}"
            shutil.copy2(ROOT/best, primary)
            # also record a normalised primary path
            sd["primary_image"] = str(primary.relative_to(ROOT))
        else:
            sd["primary_image"] = None
        slides_data.append(sd)

        # Raw text dump
        raw_text_lines.append(f"\n===== SLIDE {i}: {sd['title'] or ''} =====")
        raw_text_lines.extend(sd["texts"])
        for ti, tbl in enumerate(sd["tables"]):
            raw_text_lines.append(f"-- table {ti} --")
            for row in tbl:
                raw_text_lines.append(" | ".join(row))

        # Line-item harvest from tables on this slide
        for ti, tbl in enumerate(sd["tables"]):
            for ri, row in enumerate(tbl):
                if looks_like_lineitem_row(row):
                    raw_lineitems.append({
                        "source": "ppt",
                        "source_ref": f"slide_{i}_table_{ti}",
                        "row_index": ri,
                        "row": row,
                        "label_guess": label_from_row(row),
                        "amount_guess": amount_from_row(row),
                        "brand_guess": brand_from_row(row),
                    })

    (out_dir/"slides.json").write_text(json.dumps(slides_data, indent=2, ensure_ascii=False))
    (out_dir/"slides.txt").write_text("\n".join(raw_text_lines))
    return slides_data, raw_lineitems

def extract_docx(docx_path: Path, pkg: str):
    print(f"[docx] {pkg} <- {docx_path.name}")
    out_dir = EXTR/pkg; safe_mkdir(out_dir)
    doc = Document(str(docx_path))
    paragraphs = []
    for p in doc.paragraphs:
        text = p.text.strip()
        style = p.style.name if p.style else ""
        if not text: continue
        if "Heading" in style:
            paragraphs.append("")  # blank line before heading
            paragraphs.append(f"# {text}  [style={style}]")
        else:
            paragraphs.append(text)
    tables_out = []
    raw_lineitems = []
    for ti, t in enumerate(doc.tables):
        tbl = []
        for row in t.rows:
            tbl.append([cell.text.strip() for cell in row.cells])
        tables_out.append(tbl)
        for ri, row in enumerate(tbl):
            if looks_like_lineitem_row(row):
                raw_lineitems.append({
                    "source": "docx",
                    "source_ref": f"table_{ti}",
                    "row_index": ri,
                    "row": row,
                    "label_guess": label_from_row(row),
                    "amount_guess": amount_from_row(row),
                    "brand_guess": brand_from_row(row),
                })
    (out_dir/"docx_paragraphs.txt").write_text("\n".join(paragraphs))
    (out_dir/"docx_tables.json").write_text(json.dumps(tables_out, indent=2, ensure_ascii=False))
    return tables_out, raw_lineitems

def safe_mkdir(p):
    p = Path(p)
    # Retry against the FS-phantom race: dir may exist but exists() lies briefly
    for _ in range(8):
        try:
            if p.is_dir(): return
        except OSError:
            pass
        try:
            os.makedirs(p, exist_ok=True)
            return
        except (FileExistsError, OSError):
            import time; time.sleep(0.5)
    # last resort: trust the dir is there
    return

def main():
    safe_mkdir(EXTR)
    cross = {"by_package": {}}
    for ppt_name, pkg in PKG_MAP.items():
        ppt_path = SRC_PPT/ppt_name
        slides, ppt_items = extract_pptx(ppt_path, pkg)
        docx_name = next(d for d,p in DOCX_MAP.items() if p==pkg)
        tables, docx_items = extract_docx(SRC_DOCX/docx_name, pkg)
        all_items = ppt_items + docx_items
        # categorise
        cat_counts = {}
        for it in all_items:
            for cat in categorise(it["label_guess"] or ""):
                cat_counts[cat] = cat_counts.get(cat, 0) + 1
        # Total tables
        total_ppt_tables = sum(len(s["tables"]) for s in slides)
        cross["by_package"][pkg] = {
            "total_slides": len(slides),
            "total_ppt_tables": total_ppt_tables,
            "total_docx_tables": len(tables),
            "total_raw_lineitems": len(all_items),
            "ppt_lineitems": len(ppt_items),
            "docx_lineitems": len(docx_items),
            "category_counts": cat_counts,
        }
        # Persist raw_lineitems
        (EXTR/pkg/"raw_lineitems.json").write_text(json.dumps(all_items, indent=2, ensure_ascii=False))

    (EXTR/"cross_summary.json").write_text(json.dumps(cross, indent=2, ensure_ascii=False))
    print(json.dumps(cross, indent=2))

    # ---- Logo render test ----
    mock = ROOT/"mock"; safe_mkdir(mock)
    logo_full = ROOT/"src_logo/zuildup_logo_full.svg"
    logo_off  = ROOT/"src_logo/zuildup_logo_official.svg"
    full_b = logo_full.read_text()
    off_b  = logo_off.read_text()
    full_score = full_b.count("<path") + full_b.count("<text") + len(full_b)//200
    off_score  = off_b.count("<path")  + off_b.count("<text")  + len(off_b)//200
    # Compare: if identical content, just say so
    identical = full_b.strip() == off_b.strip()

    test_html = f"""<!doctype html><html><head><meta charset="utf-8">
<style>
body{{margin:0;background:repeating-conic-gradient(#eee 0 25%, #fff 0 50%) 0 0/40px 40px;font-family:system-ui;color:#222;}}
.row{{display:flex;gap:40px;padding:40px;align-items:flex-end}}
.cell{{background:#fff;padding:24px;border:1px solid #ddd;border-radius:12px;box-shadow:0 4px 12px rgba(0,0,0,.06)}}
.cell h3{{margin:0 0 12px;font-size:14px;font-weight:600;color:#444}}
img{{display:block}}
.dark{{background:#0A1F44;color:#fff}}
</style></head><body>
<div class="row">
  <div class="cell"><h3>logo_full @ 400px</h3><img src="../src_logo/zuildup_logo_full.svg" width="400"></div>
  <div class="cell"><h3>logo_official @ 400px</h3><img src="../src_logo/zuildup_logo_official.svg" width="400"></div>
</div>
<div class="row">
  <div class="cell"><h3>logo_full @ 800px</h3><img src="../src_logo/zuildup_logo_full.svg" width="800"></div>
</div>
<div class="row">
  <div class="cell dark"><h3>logo_full on navy @ 400px</h3><img src="../src_logo/zuildup_logo_full.svg" width="400"></div>
  <div class="cell dark"><h3>logo_official on navy @ 400px</h3><img src="../src_logo/zuildup_logo_official.svg" width="400"></div>
</div>
</body></html>"""
    (mock/"_logo_test.html").write_text(test_html)
    out_png = mock/"_logo_test.png"
    rc = os.system(f'/usr/bin/google-chrome --headless --disable-gpu --no-sandbox --hide-scrollbars '
                   f'--screenshot={out_png} --window-size=1700,1400 file://{(mock/"_logo_test.html")} 2>/tmp/chrome.log')
    rendered_ok = out_png.exists() and out_png.stat().st_size > 5000
    if rendered_ok:
        try:
            with Image.open(out_png) as im: dims = im.size
        except Exception as e:
            dims = ("err", str(e))
    else:
        dims = None
    winner = logo_full if (full_score >= off_score) else logo_off
    verdict = (
        f"# Logo verdict\n\n"
        f"- `zuildup_logo_full.svg`: {logo_full.stat().st_size}B, score={full_score} (paths/text+len heuristic)\n"
        f"- `zuildup_logo_official.svg`: {logo_off.stat().st_size}B, score={off_score}\n"
        f"- Identical content: {identical}\n"
        f"- Render PNG dims: {dims}, file: `{out_png.relative_to(ROOT)}`\n"
        f"- **Winner**: `{winner.relative_to(ROOT)}` — used as master mark for cover + spec card.\n"
        f"  Reason: {'higher path/text element score' if not identical else 'files are byte-identical, choosing zuildup_logo_full.svg by name convention'}\n"
    )
    if identical:
        winner = logo_full
    (EXTR/"logo_verdict.md").write_text(verdict)
    pkg_counts = cross["by_package"]
    print(
        f"EXTRACT_DONE "
        f"premium_slides={pkg_counts['premium']['total_slides']} "
        f"platinum_slides={pkg_counts['platinum']['total_slides']} "
        f"royale_slides={pkg_counts['royale']['total_slides']} "
        f"premium_items={pkg_counts['premium']['total_raw_lineitems']} "
        f"platinum_items={pkg_counts['platinum']['total_raw_lineitems']} "
        f"royale_items={pkg_counts['royale']['total_raw_lineitems']} "
        f"logo={winner.relative_to(ROOT)}"
    )

if __name__ == "__main__":
    try:
        main()
    except Exception:
        traceback.print_exc()
        sys.exit(1)
